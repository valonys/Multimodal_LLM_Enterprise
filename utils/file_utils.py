# utils/file_utils.py
import os
import tempfile
from pathlib import Path
from PyPDF2 import PdfReader
import pandas as pd
import docx
import pptx

UPLOAD_DIR = Path(tempfile.gettempdir()) / "uploaded_docs"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

def save_file(uploaded_file):
    dest = UPLOAD_DIR / uploaded_file.name
    with open(dest, "wb") as f:
        f.write(uploaded_file.read())
    return dest

def extract_text_from_pdf(filepath):
    reader = PdfReader(filepath)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_text_from_docx(filepath):
    doc = docx.Document(filepath)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_ppt(filepath):
    prs = pptx.Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_excel(filepath):
    dfs = pd.read_excel(filepath, sheet_name=None)
    return "\n".join([df.to_string() for df in dfs.values()])

def handle_file_upload(uploaded_files):
    results = []
    for file in uploaded_files:
        path = save_file(file)
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            results.append(extract_text_from_pdf(path))
        elif suffix in [".docx", ".doc"]:
            results.append(extract_text_from_docx(path))
        elif suffix in [".ppt", ".pptx"]:
            results.append(extract_text_from_ppt(path))
        elif suffix in [".xlsx", ".xls"]:
            results.append(extract_text_from_excel(path))
        else:
            results.append(f"Unsupported file type: {file.name}")
    return results
